# app.py
import os
import uuid
import logging
import tempfile
import base64
from datetime import datetime
from flask import Flask, request, jsonify, render_template, url_for
import chromadb
from openai import OpenAI
import fitz  # PyMuPDF for PDF processing
import pandas as pd
import docx
from pptx import Presentation
from PIL import Image
import pytesseract
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Set backend for headless operation
import io
from flask import send_file


# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("app.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Configuration
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "REMOVED_KEY")
CHROMADB_HOST = os.environ.get("CHROMADB_HOST", "chromadb-poc-0428.eastus.azurecontainer.io")
CHROMADB_PORT = int(os.environ.get("CHROMADB_PORT", "8000"))
CHROMA_COLLECTION = "document_chunks"  
TEMP_UPLOAD_FOLDER = "/tmp/temp_uploads" if os.environ.get('WEBSITE_SITE_NAME') else "temp_uploads"

# Ensure temp upload folder exists
os.makedirs(TEMP_UPLOAD_FOLDER, exist_ok=True)

# Initialize Flask app
app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # Limit uploads to 16MB

# Initialize OpenAI client
client = OpenAI(api_key=OPENAI_API_KEY)

# Connect to ChromaDB
chroma_client = chromadb.HttpClient(
    host=CHROMADB_HOST,
    port=CHROMADB_PORT,
    api_version="v2"
)

chunk_collection = chroma_client.get_or_create_collection(name=CHROMA_COLLECTION)

# Create a temporary collection for uploaded documents
uploaded_collection = chroma_client.get_or_create_collection(name="uploaded_documents")

# System prompt template for the RAG assistant in French
SYSTEM_PROMPT_FR = """Vous êtes un assistant expert en analyse de données économiques. Vos réponses sont basées sur les données économiques récupérées et vos connaissances.

Suivez ces directives pour répondre :
1. Utilisez principalement le contexte récupéré pour informer vos réponses
2. Lorsque le contexte contient des données numériques, incluez-les dans votre réponse
3. Si le contexte contient des informations contradictoires, reconnaissez-le et expliquez les différentes perspectives
4. Lorsque le contexte est insuffisant, indiquez clairement quelles informations manquent
5. Citez les sources lorsque vous faites référence à des données économiques spécifiques
6. Expliquez les concepts économiques dans un langage clair et accessible
7. N'inventez pas de données ou de statistiques qui ne figurent pas dans le contexte fourni
8. Structurez les réponses complexes avec des titres clairs et une organisation appropriée
9. Pour les données économiques sensibles au temps, mentionnez la période qu'elles couvrent

Votre objectif est de fournir une analyse précise et nuancée des données économiques pour aider les utilisateurs à prendre des décisions éclairées.
"""

# System prompt template for the RAG assistant in English
SYSTEM_PROMPT_EN = """You are an expert economic data analyst assistant. Your responses are based on retrieved economic data and your knowledge. Note that most of the source documents are in French, so you'll need to translate relevant information.

Follow these guidelines when answering:
1. Primarily use the retrieved context to inform your answers
2. When the context contains numerical data, include it in your response
3. If the context includes conflicting information, acknowledge this and explain the different perspectives
4. When the context is insufficient, clearly state what information is missing
5. Cite sources when referencing specific economic data points
6. Explain economic concepts in clear, accessible language
7. Do not make up data or statistics that aren't in the provided context
8. Structure complex responses with clear headings and organization when appropriate
9. For time-sensitive economic data, mention the time period it covers
10. Since most documents are in French, translate key information accurately

Your goal is to provide accurate, nuanced analysis of economic data to help users make informed decisions.
"""

# System prompt for uploaded document analysis
UPLOADED_DOC_PROMPT_FR = """Vous êtes un assistant expert capable d'analyser et d'extraire des informations pertinentes à partir de documents téléchargés par l'utilisateur. 

Suivez ces directives pour répondre aux questions sur le document téléchargé :
1. Utilisez principalement le contenu du document fourni pour informer vos réponses
2. Référencez des sections spécifiques du document lorsque cela est pertinent
3. Si le document contient des informations contradictoires, reconnaissez-le et expliquez les différentes perspectives
4. Pour les données numériques, assurez-vous de les citer correctement
5. Soyez précis dans vos réponses et n'inventez pas d'informations qui ne figurent pas dans le document

Votre objectif est d'aider l'utilisateur à extraire et comprendre les informations contenues dans son document.
"""

UPLOADED_DOC_PROMPT_EN = """You are an expert assistant capable of analyzing and extracting relevant information from documents uploaded by the user.

Follow these guidelines when answering questions about the uploaded document:
1. Primarily use the content of the provided document to inform your answers
2. Reference specific sections of the document when relevant
3. If the document contains conflicting information, acknowledge this and explain the different perspectives
4. For numerical data, ensure you quote it accurately
5. Be precise in your responses and do not make up information not present in the document

Your goal is to help the user extract and understand information contained in their document.
"""

def detect_language(text):
    """Simple language detection for French vs English"""
    # Common French words
    french_indicators = ['est', 'sont', 'les', 'des', 'dans', 'pour', 'avec', 'sur', 'qui', 'que', 'cette']
    
    # Check if the text contains multiple French indicators
    words = text.lower().split()
    french_word_count = sum(1 for word in words if word in french_indicators)
    
    # If more than 1 French indicators are found, consider it French
    return "fr" if french_word_count > 1 else "en"

def log_interaction(query, context, response, language, source="database"):
    """Log each interaction to help with debugging and improvement"""
    log_entry = {
        "timestamp": datetime.now().isoformat(),
        "query": query,
        "language": language,
        "source": source,
        "context_snippets": [doc[:100] + "..." for doc in context] if context else [],
        "response": response[:100] + "..."
    }
    logger.info(f"Interaction: {log_entry}")

def get_embedding(text, model="text-embedding-3-large"):
    """Generate embedding for text using OpenAI API"""
    try:
        response = client.embeddings.create(
            input=text,
            model=model
        )
        return response.data[0].embedding
    except Exception as e:
        logger.error(f"Error generating embedding: {e}")
        raise

def retrieve_context(query, top_k=5, collection=None):
    """Retrieve relevant chunks based on the query"""
    try:
        # Use the specified collection or default to the main collection
        col = collection if collection else chunk_collection
        
        # Generate embedding for query
        query_embedding = get_embedding(query)
        
        # Search ChromaDB
        results = col.query(
            query_embeddings=[query_embedding],
            n_results=top_k,
            include=["documents", "metadatas", "distances"]
        )
        
        documents = results.get("documents", [[]])[0]
        metadatas = results.get("metadatas", [[]])[0]
        distances = results.get("distances", [[]])[0]
        
        # Format context with metadata
        formatted_context = []
        for i, (doc, meta, dist) in enumerate(zip(documents, metadatas, distances)):
            source_info = f"Source: {meta.get('filename', 'Unknown')}"
            if meta.get('source_folder'):
                source_info += f" (from {meta.get('source_folder')})"
                
            formatted_doc = f"[Document {i+1}] {source_info}\n{doc}\n"
            formatted_context.append(formatted_doc)
        
        return formatted_context, documents
    except Exception as e:
        logger.error(f"Error retrieving context: {e}")
        return [], []

def generate_response(query, context_chunks, language="fr", prompt_type="default"):
    """Generate response using OpenAI model with context"""
    try:
        # Join context chunks
        joined_context = "\n\n".join(context_chunks)
        
        # Select system prompt based on language and prompt type
        if prompt_type == "uploaded":
            system_prompt = UPLOADED_DOC_PROMPT_FR if language == "fr" else UPLOADED_DOC_PROMPT_EN
        else:
            system_prompt = SYSTEM_PROMPT_FR if language == "fr" else SYSTEM_PROMPT_EN
        
        # User prompt in the appropriate language
        if language == "fr":
            user_prompt = f"""Veuillez répondre à cette question en vous basant sur le contexte suivant :

Contexte :
{joined_context}

Question : {query}
"""
        else:
            user_prompt = f"""Please help answer this question based on the following context:

Context:
{joined_context}

Question: {query}
"""
        
        # Prepare messages for chat completion
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]
        
        # Call OpenAI API
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            temperature=0.3,  # Lower temperature for more factual responses
            max_tokens=1000
        )
        
        return response.choices[0].message.content
    except Exception as e:
        logger.error(f"Error generating response: {e}")
        error_msg = "Désolé, j'ai rencontré une erreur:" if language == "fr" else "Sorry, I encountered an error:"
        return f"{error_msg} {str(e)}"

def extract_text_from_pdf(file_path):
    """Extract text from PDF file"""
    try:
        text = []
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text.append(page.get_text())
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        return f"Error extracting text from PDF: {str(e)}"

def extract_text_from_excel(file_path):
    """Extract text from Excel file"""
    try:
        df_dict = pd.read_excel(file_path, sheet_name=None)
        text_parts = []
        
        for sheet_name, df in df_dict.items():
            text_parts.append(f"Sheet: {sheet_name}")
            text_parts.append(df.to_string(index=False))
        
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting text from Excel: {e}")
        return f"Error extracting text from Excel: {str(e)}"

def extract_text_from_word(file_path):
    """Extract text from Word document"""
    try:
        doc = docx.Document(file_path)
        text = []
        
        for para in doc.paragraphs:
            text.append(para.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text.strip())
                text.append(" | ".join(row_text))
        
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Error extracting text from Word document: {e}")
        return f"Error extracting text from Word document: {str(e)}"

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint presentation"""
    try:
        prs = Presentation(file_path)
        text = []
        
        for i, slide in enumerate(prs.slides):
            text.append(f"Slide {i+1}:")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text.strip())
        
        return "\n\n".join(text)
    except Exception as e:
        logger.error(f"Error extracting text from PowerPoint: {e}")
        return f"Error extracting text from PowerPoint: {str(e)}"

def extract_text_from_image(file_path):
    """Extract text from image using OCR"""
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        logger.error(f"Error extracting text from image: {e}")
        return f"Error extracting text from image: {str(e)}"

def process_uploaded_document(file, file_type=None):
    """Process uploaded document and extract text based on file type"""
    try:
        # Create temp file path
        file_uuid = str(uuid.uuid4())
        file_extension = file.filename.split('.')[-1].lower() if '.' in file.filename else ''
        
        if not file_type:
            # Guess file type from extension if not provided
            if file_extension in ['pdf']:
                file_type = 'pdf'
            elif file_extension in ['xlsx', 'xls', 'csv']:
                file_type = 'excel'
            elif file_extension in ['docx', 'doc']:
                file_type = 'word'
            elif file_extension in ['pptx', 'ppt']:
                file_type = 'pptx'
            elif file_extension in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff']:
                file_type = 'image'
            else:
                return None, f"Unsupported file type: {file_extension}"
        
        temp_path = os.path.join(TEMP_UPLOAD_FOLDER, f"{file_uuid}.{file_extension}")
        file.save(temp_path)
        
        # Extract text based on file type
        if file_type == 'pdf':
            extracted_text = extract_text_from_pdf(temp_path)
        elif file_type == 'excel':
            extracted_text = extract_text_from_excel(temp_path)
        elif file_type == 'word':
            extracted_text = extract_text_from_word(temp_path)
        elif file_type == 'pptx':
            extracted_text = extract_text_from_pptx(temp_path)
        elif file_type == 'image':
            extracted_text = extract_text_from_image(temp_path)
        else:
            return None, f"Unsupported file type: {file_type}"
        
        # Clean up temp file
        os.remove(temp_path)
        
        # Chunk the text if it's large
        chunks = chunk_text(extracted_text, file.filename)
        
        # Add chunks to a temporary collection
        doc_collection = ingest_document_chunks(chunks, file.filename)
        
        return doc_collection, None
    except Exception as e:
        logger.error(f"Error processing uploaded document: {e}")
        return None, f"Error processing document: {str(e)}"

def chunk_text(text, filename, chunk_size=1000, overlap=200):
    """Split text into overlapping chunks"""
    chunks = []
    text_length = len(text)
    
    if text_length <= chunk_size:
        chunks.append({
            "text": text,
            "metadata": {"filename": filename, "chunk": 1, "total_chunks": 1}
        })
    else:
        # Split into overlapping chunks
        for i in range(0, text_length, chunk_size - overlap):
            chunk_text = text[i:min(i + chunk_size, text_length)]
            chunk_num = i // (chunk_size - overlap) + 1
            total_chunks = (text_length - overlap) // (chunk_size - overlap) + 1
            
            chunks.append({
                "text": chunk_text,
                "metadata": {
                    "filename": filename,
                    "chunk": chunk_num,
                    "total_chunks": total_chunks
                }
            })
    
    return chunks

def ingest_document_chunks(chunks, filename):
    """Ingest document chunks into a temporary ChromaDB collection"""
    try:
        # Create a unique collection name for this document
        collection_name = f"uploaded_{uuid.uuid4().hex}"
        doc_collection = chroma_client.get_or_create_collection(name=collection_name)
        
        # Prepare data for ingestion
        ids = []
        texts = []
        embeddings = []
        metadatas = []
        
        # Process each chunk
        for i, chunk in enumerate(chunks):
            chunk_id = f"{filename}_{i}"
            ids.append(chunk_id)
            texts.append(chunk["text"])
            
            # Generate embedding
            embedding = get_embedding(chunk["text"])
            embeddings.append(embedding)
            
            # Add metadata
            metadatas.append(chunk["metadata"])
        
        # Add data to collection
        if ids:
            doc_collection.add(
                ids=ids,
                documents=texts,
                embeddings=embeddings,
                metadatas=metadatas
            )
        
        return collection_name
    except Exception as e:
        logger.error(f"Error ingesting document chunks: {e}")
        raise



def generate_report_summary(document_collection, language="fr"):
    """Generate a summary of the uploaded document"""
    try:
        collection = chroma_client.get_collection(name=document_collection)
        
        # Get all documents from the collection
        results = collection.get(include=["documents", "metadatas"])
        
        # Concatenate all text for summary generation
        all_text = " ".join(results["documents"])
        
        # Limit text length for API efficiency
        if len(all_text) > 10000:
            all_text = all_text[:10000] + "..."
        
        # Create a summary prompt
        if language == "fr":
            summary_prompt = f"""Voici un document dont j'ai besoin d'un résumé concis et structuré. 
            Merci de fournir les sections suivantes:
            1. Points clés (5 points maximum)
            2. Résumé exécutif (maximum 3 paragraphes)
            3. Données importantes (chiffres, statistiques, dates, etc.)
            
            Document:
            {all_text}
            """
        else:
            summary_prompt = f"""Here is a document that I need a concise and structured summary for.
            Please provide the following sections:
            1. Key points (maximum 5 points)
            2. Executive summary (maximum 3 paragraphs)
            3. Important data (figures, statistics, dates, etc.)
            
            Document:
            {all_text}
            """
        
        # Generate summary using OpenAI
        messages = [
            {"role": "system", "content": "You are an expert in document summarization."},
            {"role": "user", "content": summary_prompt}
        ]
        
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            temperature=0.3,
            max_tokens=800
        )
        
        return response.choices[0].message.content
    except Exception as e:
        logger.error(f"Error generating report summary: {e}")
        error_msg = "Désolé, je n'ai pas pu générer le résumé:" if language == "fr" else "Sorry, I couldn't generate the summary:"
        return f"{error_msg} {str(e)}"



def extract_data_for_visualization(document_collection, language="fr"):
    """Extract data suitable for visualization from the document with enhanced capabilities"""
    try:
        collection = chroma_client.get_collection(name=document_collection)
        
        # Get all documents from the collection
        results = collection.get(include=["documents", "metadatas"])
        
        # Concatenate all text for analysis
        all_text = " ".join(results["documents"])
        
        # Limit text length for API efficiency
        if len(all_text) > 12000:
            all_text = all_text[:12000] + "..."
        
        # Enhanced data extraction prompt with more chart types and structure
        if language == "fr":
            data_prompt = f"""Veuillez extraire des données numériques qui pourraient être visualisées à partir du document suivant.
            Pour chaque jeu de données potentiel, spécifiez:
            1. Le type de visualisation le plus approprié parmi les options suivantes:
               - "bar" (graphique à barres) - pour comparer des valeurs entre catégories
               - "line" (graphique linéaire) - pour montrer l'évolution dans le temps
               - "pie" (camembert) - pour montrer des proportions d'un tout
               - "scatter" (nuage de points) - pour montrer des corrélations entre deux variables
               - "heatmap" (carte thermique) - pour visualiser l'intensité de données sur une matrice
               - "radar" (graphique radar) - pour comparer plusieurs variables quantitatives
               - "bubble" (graphique à bulles) - pour comparer trois dimensions de données
               - "treemap" (carte proportionnelle) - pour visualiser des données hiérarchiques
               - "area" (graphique de zones) - pour l'évolution de volumes dans le temps
            
            2. Les points de données exacts (maximum 10 points)
            3. Un titre significatif pour cette visualisation
            4. Une courte description (1-2 phrases) expliquant l'importance de ces données
            
            Pour les graphiques complexes comme les heatmaps, fournissez les données sous forme de matrice.
            
            Fournissez les données dans un format structuré et lisible par machine.
            Exemple:
            [
              {{
                "type": "bar",
                "title": "Ventes par région",
                "description": "Montre la distribution des ventes à travers différentes régions pour l'année 2023.",
                "labels": ["Nord", "Sud", "Est", "Ouest"],
                "values": [123, 456, 789, 321],
                "x_axis_label": "Région",
                "y_axis_label": "Ventes (€)"
              }},
              {{
                "type": "heatmap",
                "title": "Matrice de corrélation des indicateurs économiques",
                "description": "Montre les corrélations entre différents indicateurs économiques.",
                "labels": ["PIB", "Inflation", "Chômage", "Exportations"],
                "values": [
                  [1.0, 0.2, -0.5, 0.7],
                  [0.2, 1.0, 0.1, -0.3],
                  [-0.5, 0.1, 1.0, -0.6],
                  [0.7, -0.3, -0.6, 1.0]
                ]
              }}
            ]
            
            Si le document contient des données tabulaires, veuillez également les extraire dans un format adapté pour un tableau dans le rapport.
            
            Document:
            {all_text}
            """
        else:
            data_prompt = f"""Please extract numerical data that could be visualized from the following document.
            For each potential dataset, specify:
            1. The most appropriate visualization type from the following options:
               - "bar" (bar chart) - for comparing values across categories
               - "line" (line graph) - for showing evolution over time
               - "pie" (pie chart) - for showing proportions of a whole
               - "scatter" (scatter plot) - for showing correlations between two variables
               - "heatmap" (heat map) - for visualizing intensity of data on a matrix
               - "radar" (radar chart) - for comparing multiple quantitative variables
               - "bubble" (bubble chart) - for comparing three dimensions of data
               - "treemap" (treemap) - for visualizing hierarchical data
               - "area" (area chart) - for evolution of volumes over time
            
            2. The exact data points (maximum 10 points)
            3. A meaningful title for this visualization
            4. A short description (1-2 sentences) explaining the importance of this data
            
            For complex charts like heatmaps, provide the data as a matrix.
            
            Provide the data in a structured, machine-readable format.
            Example:
            [
              {{
                "type": "bar",
                "title": "Sales by Region",
                "description": "Shows the distribution of sales across different regions for the year 2023.",
                "labels": ["North", "South", "East", "West"],
                "values": [123, 456, 789, 321],
                "x_axis_label": "Region",
                "y_axis_label": "Sales ($)"
              }},
              {{
                "type": "heatmap",
                "title": "Correlation Matrix of Economic Indicators",
                "description": "Shows the correlations between different economic indicators.",
                "labels": ["GDP", "Inflation", "Unemployment", "Exports"],
                "values": [
                  [1.0, 0.2, -0.5, 0.7],
                  [0.2, 1.0, 0.1, -0.3],
                  [-0.5, 0.1, 1.0, -0.6],
                  [0.7, -0.3, -0.6, 1.0]
                ]
              }}
            ]
            
            If the document contains tabular data, please also extract it in a format suitable for a table in the report.
            
            Document:
            {all_text}
            """
        
        # Generate data extraction using OpenAI
        messages = [
            {"role": "system", "content": "You are an expert in data extraction and visualization."},
            {"role": "user", "content": data_prompt}
        ]
        
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            temperature=0.3,
            max_tokens=1500
        )
        
        # Try to parse the response as JSON
        data_str = response.choices[0].message.content
        
        # Attempt to find and extract JSON content
        import re
        import json
        
        # Look for JSON-like content with brackets
        json_match = re.search(r'\[\s*\{.*\}\s*\]', data_str, re.DOTALL)
        if json_match:
            try:
                data = json.loads(json_match.group(0))
                return data
            except json.JSONDecodeError as e:
                logger.error(f"Error parsing JSON from response: {e}")
        
        # If we can't parse JSON, return an empty list
        return []
    except Exception as e:
        logger.error(f"Error extracting data for visualization: {e}")
        return []




def extract_table_data(document_collection, language="fr"):
    """Extract tabular data from the document for report tables"""
    try:
        collection = chroma_client.get_collection(name=document_collection)
        
        # Get all documents from the collection
        results = collection.get(include=["documents", "metadatas"])
        
        # Concatenate all text for analysis
        all_text = " ".join(results["documents"])
        
        # Limit text length for API efficiency
        if len(all_text) > 10000:
            all_text = all_text[:10000] + "..."
        
        # Table extraction prompt
        if language == "fr":
            table_prompt = f"""Veuillez extraire des données tabulaires à partir du document suivant.
            Pour chaque tableau potentiel, spécifiez:
            1. Un titre descriptif pour le tableau
            2. Les en-têtes de colonnes
            3. Les données pour chaque ligne
            4. Une brève description (1-2 phrases) expliquant l'importance de ces données
            
            Fournissez les données dans un format structuré et lisible par machine.
            Exemple:
            [
              {{
                "title": "Indicateurs économiques par pays",
                "description": "Comparaison des indicateurs économiques clés entre les principaux pays européens.",
                "headers": ["Pays", "PIB (M€)", "Inflation (%)", "Chômage (%)"],
                "data": [
                  ["France", "2 500 000", "2.1", "7.8"],
                  ["Allemagne", "3 800 000", "1.8", "3.4"],
                  ["Italie", "1 900 000", "2.7", "9.1"],
                  ["Espagne", "1 400 000", "3.2", "14.2"]
                ]
              }}
            ]
            
            Document:
            {all_text}
            """
        else:
            table_prompt = f"""Please extract tabular data from the following document.
            For each potential table, specify:
            1. A descriptive title for the table
            2. The column headers
            3. The data for each row
            4. A brief description (1-2 sentences) explaining the importance of this data
            
            Provide the data in a structured, machine-readable format.
            Example:
            [
              {{
                "title": "Economic Indicators by Country",
                "description": "Comparison of key economic indicators across major European countries.",
                "headers": ["Country", "GDP (M€)", "Inflation (%)", "Unemployment (%)"],
                "data": [
                  ["France", "2,500,000", "2.1", "7.8"],
                  ["Germany", "3,800,000", "1.8", "3.4"],
                  ["Italy", "1,900,000", "2.7", "9.1"],
                  ["Spain", "1,400,000", "3.2", "14.2"]
                ]
              }}
            ]
            
            Document:
            {all_text}
            """
        
        # Generate table extraction using OpenAI
        messages = [
            {"role": "system", "content": "You are an expert in data extraction and table formatting."},
            {"role": "user", "content": table_prompt}
        ]
        
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            temperature=0.3,
            max_tokens=1200
        )
        
        # Try to parse the response as JSON
        data_str = response.choices[0].message.content
        
        # Attempt to find and extract JSON content
        import re
        import json
        
        # Look for JSON-like content with brackets
        json_match = re.search(r'\[\s*\{.*\}\s*\]', data_str, re.DOTALL)
        if json_match:
            try:
                data = json.loads(json_match.group(0))
                return data
            except json.JSONDecodeError as e:
                logger.error(f"Error parsing JSON from response: {e}")
        
        # If we can't parse JSON, return an empty list
        return []
    except Exception as e:
        logger.error(f"Error extracting table data: {e}")
        return []



def create_visualizations(data):
    """Create enhanced visualizations based on extracted data"""
    import matplotlib.pyplot as plt
    import matplotlib.cm as cm
    import numpy as np
    import seaborn as sns
    from matplotlib.colors import LinearSegmentedColormap
    
    # Set modern style
    plt.style.use('seaborn-v0_8-whitegrid')
    
    # Modern color palette
    modern_colors = ['#3498db', '#2ecc71', '#e74c3c', '#f39c12', '#9b59b6', 
                    '#1abc9c', '#d35400', '#34495e', '#7f8c8d', '#2c3e50']
    
    visualizations = []
    
    for i, chart_data in enumerate(data):
        try:
            # Extract data
            chart_type = chart_data.get("type", "bar").lower()
            title = chart_data.get("title", f"Chart {i+1}")
            description = chart_data.get("description", "")
            labels = chart_data.get("labels", [])
            values = chart_data.get("values", [])
            x_axis = chart_data.get("x_axis_label", "")
            y_axis = chart_data.get("y_axis_label", "")
            
            # Create figure with appropriate size for the chart type
            if chart_type in ["heatmap", "radar", "treemap"]:
                plt.figure(figsize=(10, 8))
            else:
                plt.figure(figsize=(10, 6))
            
            # Create different chart types
            if chart_type == "pie":
                # Enhanced pie chart
                plt.pie(values, labels=None, autopct='%1.1f%%', startangle=90, 
                       colors=modern_colors[:len(values)], shadow=False)
                plt.title(title, fontsize=16, pad=20)
                # Add legend outside the pie
                plt.legend(labels, loc="center left", bbox_to_anchor=(1, 0, 0.5, 1))
                
            elif chart_type == "line":
                # Enhanced line chart
                plt.plot(range(len(labels)), values, marker='o', linewidth=3, 
                        color=modern_colors[0])
                plt.title(title, fontsize=16, pad=20)
                plt.xlabel(x_axis if x_axis else "Categories", fontsize=12)
                plt.ylabel(y_axis if y_axis else "Values", fontsize=12)
                plt.xticks(range(len(labels)), labels, rotation=45)
                plt.grid(True, linestyle='--', alpha=0.7)
                
            elif chart_type == "bar":
                # Enhanced bar chart
                bars = plt.bar(range(len(labels)), values, color=modern_colors[:len(values)])
                plt.title(title, fontsize=16, pad=20)
                plt.xlabel(x_axis if x_axis else "Categories", fontsize=12)
                plt.ylabel(y_axis if y_axis else "Values", fontsize=12)
                plt.xticks(range(len(labels)), labels, rotation=45)
                # Add value labels on top of bars
                for bar in bars:
                    height = bar.get_height()
                    plt.text(bar.get_x() + bar.get_width()/2., height + 0.01,
                           f'{height:,.1f}', ha='center', fontsize=10)
                
            elif chart_type == "scatter":
                # Scatter plot (assumes two sets of values)
                if isinstance(values[0], list) and len(values) >= 2:
                    x_values = values[0]
                    y_values = values[1]
                    plt.scatter(x_values, y_values, s=100, alpha=0.7, c=modern_colors[0])
                    plt.title(title, fontsize=16, pad=20)
                    plt.xlabel(x_axis if x_axis else labels[0], fontsize=12)
                    plt.ylabel(y_axis if y_axis else labels[1], fontsize=12)
                    plt.grid(True, linestyle='--', alpha=0.7)
                else:
                    # Fallback to scatter against index
                    plt.scatter(range(len(values)), values, s=100, alpha=0.7, c=modern_colors[0])
                    plt.title(title, fontsize=16, pad=20)
                    plt.xlabel(x_axis if x_axis else "Index", fontsize=12)
                    plt.ylabel(y_axis if y_axis else "Values", fontsize=12)
                    plt.xticks(range(len(labels)), labels, rotation=45)
                    plt.grid(True, linestyle='--', alpha=0.7)
                
            elif chart_type == "heatmap":
                # Heatmap visualization
                if isinstance(values[0], list):
                    # Create a custom colormap
                    colors = ["#f1c40f", "#e67e22", "#e74c3c", "#9b59b6", "#3498db"]
                    n_bins = 100
                    cmap_name = 'custom_cmap'
                    cm = LinearSegmentedColormap.from_list(cmap_name, colors, N=n_bins)
                    
                    # Create the heatmap
                    sns.heatmap(values, annot=True, cmap=cm, fmt=".2f", 
                               xticklabels=labels, yticklabels=labels)
                    plt.title(title, fontsize=16, pad=20)
                else:
                    # Fallback to bar chart if data isn't suitable for heatmap
                    plt.bar(range(len(labels)), values, color=modern_colors[:len(values)])
                    plt.title(title + " (Fallback)", fontsize=16, pad=20)
                    plt.xlabel("Categories", fontsize=12)
                    plt.ylabel("Values", fontsize=12)
                    plt.xticks(range(len(labels)), labels, rotation=45)
                
            elif chart_type == "radar":
                # Radar chart (polar plot)
                # Convert to numpy arrays for easier manipulation
                values_np = np.array(values)
                
                # Number of variables
                N = len(labels)
                
                # Compute angle for each axis
                angles = [n / float(N) * 2 * np.pi for n in range(N)]
                angles += angles[:1]  # Close the loop
                
                # Add the last value to close the polygon
                values_radar = np.concatenate((values_np, [values_np[0]]))
                
                # Draw the plot
                ax = plt.subplot(111, polar=True)
                ax.plot(angles, values_radar, linewidth=2, linestyle='solid', color=modern_colors[0])
                ax.fill(angles, values_radar, alpha=0.25, color=modern_colors[0])
                
                # Set labels
                ax.set_xticks(angles[:-1])
                ax.set_xticklabels(labels)
                
                # Add title
                plt.title(title, fontsize=16, pad=20)
                
            elif chart_type == "bubble":
                # Bubble chart (scatter with sized bubbles)
                # This assumes we have at least 3 dimensions: x, y, and size
                if isinstance(values[0], list) and len(values) >= 3:
                    x_values = values[0]
                    y_values = values[1]
                    sizes = [s * 100 for s in values[2]]  # Scale up sizes for visibility
                    
                    plt.scatter(x_values, y_values, s=sizes, alpha=0.6, c=modern_colors[0])
                    plt.title(title, fontsize=16, pad=20)
                    plt.xlabel(x_axis if x_axis else labels[0], fontsize=12)
                    plt.ylabel(y_axis if y_axis else labels[1], fontsize=12)
                    plt.grid(True, linestyle='--', alpha=0.7)
                else:
                    # Fallback to regular scatter
                    plt.scatter(range(len(values)), values, s=100, alpha=0.7, c=modern_colors[0])
                    plt.title(title + " (Fallback)", fontsize=16, pad=20)
                    plt.xlabel("Categories", fontsize=12)
                    plt.ylabel("Values", fontsize=12)
                    plt.xticks(range(len(labels)), labels, rotation=45)
                    plt.grid(True, linestyle='--', alpha=0.7)
                
            elif chart_type == "area":
                # Area chart
                plt.fill_between(range(len(labels)), values, alpha=0.5, color=modern_colors[0])
                plt.plot(range(len(labels)), values, color=modern_colors[0], alpha=0.8)
                plt.title(title, fontsize=16, pad=20)
                plt.xlabel(x_axis if x_axis else "Categories", fontsize=12)
                plt.ylabel(y_axis if y_axis else "Values", fontsize=12)
                plt.xticks(range(len(labels)), labels, rotation=45)
                plt.grid(True, linestyle='--', alpha=0.7)
                
            elif chart_type == "treemap":
                # Simple treemap implementation using nested rectangles
                # This is a simplified version; for complex treemaps, consider using squarify
                try:
                    import squarify
                    
                    # Normalize values for treemap
                    norm_values = [v/sum(values) * 1000 for v in values]
                    
                    # Create treemap
                    squarify.plot(sizes=norm_values, label=labels, color=modern_colors[:len(values)],
                                 alpha=0.8, pad=0.5)
                    plt.axis('off')
                    plt.title(title, fontsize=16, pad=20)
                except ImportError:
                    # Fallback to pie chart if squarify is not available
                    plt.pie(values, labels=labels, autopct='%1.1f%%', startangle=90,
                          colors=modern_colors[:len(values)])
                    plt.title(title + " (Fallback as Pie)", fontsize=16, pad=20)
            
            else:  # Default to bar chart
                plt.bar(range(len(labels)), values, color=modern_colors[:len(values)])
                plt.title(title, fontsize=16, pad=20)
                plt.xlabel("Categories", fontsize=12)
                plt.ylabel("Values", fontsize=12)
                plt.xticks(range(len(labels)), labels, rotation=45)
            
            # Add description as a text box if available
            if description:
                plt.figtext(0.5, 0.01, description, wrap=True, horizontalalignment='center', 
                          fontsize=10, style='italic')
            
            plt.tight_layout()
            
            # Save the figure to a BytesIO object
            img_data = io.BytesIO()
            plt.savefig(img_data, format='png', dpi=150, bbox_inches='tight')
            img_data.seek(0)
            
            # Add the image data to the visualizations list
            visualizations.append({
                "title": title,
                "description": description,
                "image_data": img_data,
                "chart_type": chart_type
            })
            
            # Close the figure to free memory
            plt.close()
        except Exception as e:
            logger.error(f"Error creating visualization: {e}")
    
    return visualizations



def generate_pdf_report(summary, visualizations, tables, filename, language="fr"):
    """Generate an enhanced PDF report with summary, visualizations, and tables"""
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Image, 
                                   Table, TableStyle, PageBreak, ListFlowable, 
                                   ListItem, Flowable, KeepTogether)
    import io
    from datetime import datetime
    
    try:
        # Create a BytesIO object to store the PDF
        pdf_buffer = io.BytesIO()
        
        # Create the PDF document with better margins
        doc = SimpleDocTemplate(
            pdf_buffer, 
            pagesize=letter,
            leftMargin=0.75*inch,
            rightMargin=0.75*inch,
            topMargin=0.75*inch,
            bottomMargin=0.75*inch
        )
        
        # Enhanced styles
        styles = getSampleStyleSheet()
        
        # Custom styles for better hierarchy
        title_style = ParagraphStyle(
            'ReportTitle',
            parent=styles['Heading1'],
            fontSize=20,
            spaceAfter=24,
            alignment=1,  # Center alignment
            textColor=colors.HexColor('#2c3e50')
        )
        
        section_style = ParagraphStyle(
            'SectionHeader',
            parent=styles['Heading1'],
            fontSize=16,
            spaceBefore=16,
            spaceAfter=12,
            textColor=colors.HexColor('#2980b9'),
            borderWidth=0,
            borderPadding=0,
            borderColor=colors.HexColor('#bdc3c7'),
            borderRadius=None
        )
        
        subsection_style = ParagraphStyle(
            'SubSectionHeader',
            parent=styles['Heading2'],
            fontSize=14,
            spaceBefore=12,
            spaceAfter=8,
            textColor=colors.HexColor('#3498db')
        )
        
        normal_style = ParagraphStyle(
            'BodyText',
            parent=styles['Normal'],
            fontSize=11,
            leading=14,
            spaceBefore=6,
            spaceAfter=6
        )
        
        caption_style = ParagraphStyle(
            'Caption',
            parent=styles['Italic'],
            fontSize=10,
            leading=12,
            alignment=1  # Center alignment
        )
        
        # Custom Flowable for horizontal lines
        class HorizontalLine(Flowable):
            def __init__(self, width, color=colors.HexColor('#ecf0f1'), thickness=1):
                Flowable.__init__(self)
                self.width = width
                self.color = color
                self.thickness = thickness
            
            def draw(self):
                self.canv.setStrokeColor(self.color)
                self.canv.setLineWidth(self.thickness)
                self.canv.line(0, 0, self.width, 0)
        
        # Elements to add to the PDF
        elements = []
        
        # Add title
        if language == "fr":
            title = "Rapport d'Analyse de Document"
            summary_title = "Résumé du Document"
            viz_title = "Visualisations"
            tables_title = "Tableaux de Données"
            gen_on = "Généré le"
        else:
            title = "Document Analysis Report"
            summary_title = "Document Summary"
            viz_title = "Visualizations"
            tables_title = "Data Tables"
            gen_on = "Generated on"
            
        elements.append(Paragraph(title, title_style))
        elements.append(HorizontalLine(450, colors.HexColor('#bdc3c7'), 2))
        elements.append(Spacer(1, 12))
        
        # Add summary with proper formatting
        elements.append(Paragraph(summary_title, section_style))
        
        # Process the summary to extract structure if possible
        try:
            # Try to identify sections with bullet points or numbered lists
            import re
            
            # Look for sections with potential headings
            section_pattern = r'(.*?)\n((?:[\s\S](?!\n\n[A-Z0-9]))*)'
            sections = re.findall(section_pattern, summary + "\n\n")
            
            if sections:
                for heading, content in sections:
                    if heading.strip():
                        # Check if this looks like a heading
                        if re.match(r'^[A-Z0-9].*[:.)]$', heading.strip()):
                            # It's likely a heading
                            elements.append(Paragraph(heading.strip(), subsection_style))
                        else:
                            # It's regular content
                            elements.append(Paragraph(heading.strip(), normal_style))
                    
                    # Process content which might have bullet points
                    if content.strip():
                        # Check for bullet points or numbered lists
                        bullet_pattern = r'^\s*[\*\-•]\s+(.*)'
                        numbered_pattern = r'^\s*(\d+)[\.)]?\s+(.*)'
                        
                        lines = content.strip().split('\n')
                        current_list = []
                        in_list = False
                        
                        for line in lines:
                            bullet_match = re.match(bullet_pattern, line)
                            numbered_match = re.match(numbered_pattern, line)
                            
                            if bullet_match or numbered_match:
                                if not in_list:
                                    in_list = True
                                    current_list = []
                                
                                list_text = bullet_match.group(1) if bullet_match else numbered_match.group(2)
                                current_list.append(ListItem(Paragraph(list_text, normal_style)))
                            else:
                                # Not a list item
                                if in_list:
                                    # End previous list
                                    elements.append(ListFlowable(current_list, bulletType='bullet'))
                                    in_list = False
                                
                                if line.strip():
                                    elements.append(Paragraph(line.strip(), normal_style))
                        
                        # End any remaining list
                        if in_list:
                            elements.append(ListFlowable(current_list, bulletType='bullet'))
            else:
                # Fallback to simple paragraph splitting if no structure detected
                for paragraph in summary.split('\n\n'):
                    if paragraph.strip():
                        elements.append(Paragraph(paragraph.strip(), normal_style))
        except Exception as e:
            # Fallback if structured parsing fails
            logger.error(f"Error parsing summary structure: {e}")
            for paragraph in summary.split('\n\n'):
                if paragraph.strip():
                    elements.append(Paragraph(paragraph.strip(), normal_style))
        
        # Add separator
        elements.append(Spacer(1, 12))
        elements.append(HorizontalLine(450, colors.HexColor('#bdc3c7'), 1))
        elements.append(Spacer(1, 12))
        
        # Add tables if available
        if tables:
            elements.append(Paragraph(tables_title, section_style))
            
            for table_data in tables:
                # Add table title
                table_title = table_data.get("title", "Data Table")
                elements.append(Paragraph(table_title, subsection_style))
                
                # Add table description if available
                description = table_data.get("description", "")
                if description:
                    elements.append(Paragraph(description, normal_style))
                    elements.append(Spacer(1, 6))
                
                # Create table
                headers = table_data.get("headers", [])
                data = table_data.get("data", [])
                
                if headers and data:
                    # Create table data with headers
                    table_content = [headers]
                    table_content.extend(data)
                    
                    # Create ReportLab table
                    table = Table(table_content, repeatRows=1)
                    
                    # Add table styling
                    style = TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#3498db')),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 10),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#f9f9f9')),
                        ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                        ('ALIGN', (0, 1), (-1, -1), 'LEFT'),
                        ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                        ('FONTSIZE', (0, 1), (-1, -1), 9),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#bdc3c7')),
                        ('BOX', (0, 0), (-1, -1), 0.5, colors.HexColor('#7f8c8d')),
                        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.HexColor('#f9f9f9'), colors.HexColor('#ffffff')])
                    ])
                    
                    # Apply alternating row colors
                    for i in range(1, len(table_content)):
                        if i % 2 == 0:
                            style.add('BACKGROUND', (0, i), (-1, i), colors.HexColor('#f2f2f2'))
                    
                    table.setStyle(style)
                    
                    # Make the table a KeepTogether to avoid breaking across pages if possible
                    elements.append(KeepTogether([table, Spacer(1, 6)]))
                else:
                    elements.append(Paragraph("No data available for this table.", caption_style))
                
                elements.append(Spacer(1, 12))
        
        # Add page break before visualizations if we have both tables and visualizations
        if tables and visualizations:
            elements.append(PageBreak())
        
        # Add visualizations
        if visualizations:
            elements.append(Paragraph(viz_title, section_style))
            
            for i, viz in enumerate(visualizations):
                # Add new page for each visualization after the first one
                if i > 0:
                    elements.append(PageBreak())
                
                # Extract visualization data
                title = viz.get("title", f"Chart {i+1}")
                description = viz.get("description", "")
                chart_type = viz.get("chart_type", "")
                img_data = viz.get("image_data")
                
                # Add chart title with chart type
                chart_type_text = f" ({chart_type.capitalize()})" if chart_type else ""
                elements.append(Paragraph(f"{title}{chart_type_text}", subsection_style))
                
                # Add chart description if available
                if description:
                    elements.append(Paragraph(description, normal_style))
                    elements.append(Spacer(1, 6))
                
                # Add chart image with responsive sizing
                if img_data:
                    img = Image(img_data, width=6*inch, height=4.5*inch)
                    elements.append(img)
                    elements.append(Spacer(1, 12))
            
        # Add timestamp at the end
        elements.append(Spacer(1, 24))
        elements.append(HorizontalLine(450, colors.HexColor('#bdc3c7'), 1))
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        elements.append(Paragraph(f"{gen_on} {timestamp}", caption_style))
        
        # Build the PDF
        doc.build(elements)
        
        # Get the PDF data
        pdf_data = pdf_buffer.getvalue()
        pdf_buffer.close()
        
        return pdf_data
        
    except Exception as e:
        logger.error(f"Error generating PDF report: {e}")
        raise






@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/query', methods=['POST'])
def process_query():
    data = request.json
    query = data.get('query', '')
    
    if not query:
        return jsonify({"error": "Query is required"}), 400
    
    # Detect language or use provided language
    language = data.get('language', detect_language(query))
    
    # Check if there's a document collection to use
    doc_collection = data.get('document_collection')
    
    if doc_collection:
        # Try to get the collection
        try:
            collection = chroma_client.get_collection(name=doc_collection)
            context_chunks, raw_docs = retrieve_context(query, collection=collection)
            response = generate_response(query, context_chunks, language, prompt_type="uploaded")
            source = "uploaded_document"
        except Exception as e:
            logger.error(f"Error using document collection: {e}")
            # Fall back to the main collection
            context_chunks, raw_docs = retrieve_context(query)
            response = generate_response(query, context_chunks, language)
            source = "database"
    else:
        # Use the main collection
        context_chunks, raw_docs = retrieve_context(query)
        response = generate_response(query, context_chunks, language)
        source = "database"
    
    # Log interaction
    log_interaction(query, raw_docs, response, language, source=source)
    
    return jsonify({
        "query": query,
        "response": response,
        "context": context_chunks[:3],  # Return first 3 context chunks for reference
        "language": language,
        "document_collection": doc_collection
    })

@app.route('/api/upload', methods=['POST'])
def upload_document():
    # Check if the post request has the file part
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    # Process the uploaded document
    doc_collection, error = process_uploaded_document(file)
    
    if error:
        return jsonify({"error": error}), 400
    
    return jsonify({
        "success": True,
        "message": "Document uploaded and processed successfully",
        "document_collection": doc_collection,
        "filename": file.filename
    })

@app.route('/health')
def health_check():
    return jsonify({"status": "ok"})



@app.route('/api/generate-report', methods=['POST'])
def generate_report():
    """Generate a report for an uploaded document"""
    data = request.json
    document_collection = data.get('document_collection')
    filename = data.get('filename', 'document')
    language = data.get('language', DEFAULT_LANGUAGE)
    
    if not document_collection:
        return jsonify({"error": "Document collection ID is required"}), 400
    
    try:
        # Generate summary
        summary = generate_report_summary(document_collection, language)
        
        # Extract data for visualizations
        data_for_viz = extract_data_for_visualization(document_collection, language)
        
        # Extract table data
        tables_data = extract_table_data(document_collection, language)
        
        # Create visualizations
        visualizations = create_visualizations(data_for_viz)
        
        # Generate PDF
        pdf_data = generate_pdf_report(summary, visualizations, tables_data, filename, language)
        
        if not pdf_data:
            logger.error("PDF data is None or empty")
            return jsonify({"error": "Generated report is empty"}), 500
            
        # Encode PDF data as base64 for transmission
        pdf_base64 = base64.b64encode(pdf_data).decode('utf-8')
        
        return jsonify({
            "success": True,
            "message": "Report generated successfully",
            "report_data": pdf_base64,
            "filename": f"{filename}_report.pdf"
        })
    except Exception as e:
        logger.error(f"Error in report generation: {e}", exc_info=True)  # Added exc_info for stack trace
        return jsonify({"error": f"Failed to generate report: {str(e)}"}), 500


        
@app.route('/api/download-report/<filename>', methods=['GET'])
def download_report(filename):
    """Download a generated report"""
    # This route would handle downloading a saved report
    # For now, we'll just return an error as reports are generated on-demand
    return jsonify({"error": "Direct report downloads not yet implemented"}), 501





if __name__ == '__main__':
    # Create templates directory if it doesn't exist
    if not os.path.exists('templates'):
        os.makedirs('templates')
    
    # Run the app
    app.run(host='0.0.0.0', port=PORT, debug=True)